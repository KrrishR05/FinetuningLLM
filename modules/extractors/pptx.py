from __future__ import annotations

import io
from typing import List, Tuple

from modules.cleaner import full_clean
from modules.extractors._base import BaseExtractor
from modules.models import PageContent, Section, TableData


class PPTXExtractor(BaseExtractor):

    def extract(self, data: bytes) -> Tuple[List[PageContent], List[str]]:
        from pptx import Presentation

        pages: List[PageContent] = []
        warnings: List[str] = []

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as e:
            warnings.append(f"Failed to open PPTX: {e}")
            return pages, warnings

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            texts: list[str] = []
            sections: list[Section] = []
            tables: list[TableData] = []

            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title_text = slide.shapes.title.text.strip()
                    sections.append(Section(
                        title=title_text,
                        level=1,
                        text="",
                        page_number=slide_num,
                    ))

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                texts.append(para_text)

                    if shape.has_table:
                        table = shape.table
                        t_rows = []
                        for row in table.rows:
                            row_texts = [cell.text.strip() for cell in row.cells]
                            if any(row_texts):
                                t_rows.append(row_texts)
                                texts.append(" | ".join(t for t in row_texts if t))

                        if t_rows:
                            headers = t_rows[0]
                            data_rows = t_rows[1:] if len(t_rows) > 1 else []
                            tables.append(TableData(
                                headers=headers,
                                rows=data_rows,
                                page_number=slide_num,
                            ))
            except Exception as e:
                warnings.append(f"Error reading Slide {slide_num}: {e}")

            text = full_clean("\n".join(texts))
            if not text.strip():
                warnings.append(f"Slide {slide_num} had no extractable text")

            pages.append(PageContent(
                page_number=slide_num,
                text=text,
                source_label=f"Slide {slide_num}",
                sections=sections,
                tables=tables,
            ))

        return pages, warnings


def extract_pptx(data: bytes) -> Tuple[List[PageContent], List[str]]:
    return PPTXExtractor().extract(data)
